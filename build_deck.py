# -*- coding: utf-8 -*-
"""Build the CV class presentation for the Korean card-news layout detector."""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_LABEL_POSITION

ROOT = os.path.dirname(os.path.abspath(__file__))
def P(*a): return os.path.join(ROOT, *a)

# ---- palette ----
NAVY   = RGBColor(0x14, 0x18, 0x2E)   # deep background
PANEL  = RGBColor(0x1F, 0x25, 0x45)   # dark panel
INK    = RGBColor(0x1A, 0x1F, 0x3D)   # dark text on light
LIGHT  = RGBColor(0xF4, 0xF6, 0xFC)   # light bg / light text
MUTE_D = RGBColor(0xA9, 0xB0, 0xCC)   # muted on dark
MUTE_L = RGBColor(0x60, 0x6A, 0x8C)   # muted on light
CORAL  = RGBColor(0xFF, 0x5C, 0x7A)   # accent (energy)
MINT   = RGBColor(0x16, 0xC0, 0x8E)   # metric / success
AMBER  = RGBColor(0xF5, 0xA6, 0x23)   # secondary accent
CARDBG = RGBColor(0xFF, 0xFF, 0xFF)   # white card on light
TINT   = RGBColor(0xEE, 0xF1, 0xFB)   # subtle tint card

FONT = "Malgun Gothic"

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
SW, SH = 13.333, 7.5

def slide(bg=LIGHT):
    s = prs.slides.add_slide(BLANK)
    r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    r.fill.solid(); r.fill.fore_color.rgb = bg; r.line.fill.background()
    r.shadow.inherit = False
    s.shapes._spTree.remove(r._element); s.shapes._spTree.insert(2, r._element)
    return s

def box(s, x, y, w, h, fill=None, line=None, lw=1.0, radius=False, shadow=False):
    shp = s.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE,
        Inches(x), Inches(y), Inches(w), Inches(h))
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line; shp.line.width = Pt(lw)
    shp.shadow.inherit = False
    if shadow:
        from pptx.oxml.ns import qn
        sp = shp.shadow._element  # this IS the spPr (CT_ShapeProperties)
        for ex in sp.findall(qn('a:effectLst')):
            sp.remove(ex)
        ef = sp.makeelement(qn('a:effectLst'), {})
        sh = sp.makeelement(qn('a:outerShdw'),
                            {'blurRad':'90000','dist':'38100','dir':'5400000','rotWithShape':'0'})
        clr = sp.makeelement(qn('a:srgbClr'), {'val':'1A1F3D'})
        al = sp.makeelement(qn('a:alpha'), {'val':'18000'})
        clr.append(al); sh.append(clr); ef.append(sh); sp.append(ef)
    return shp

def text(s, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
         space_after=4, line_spacing=1.0):
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    if isinstance(runs[0], tuple): runs = [runs]
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align; p.space_after = Pt(space_after); p.space_before = Pt(0)
        p.line_spacing = line_spacing
        for (t, sz, col, bold) in para:
            r = p.add_run(); r.text = t
            r.font.name = FONT; r.font.size = Pt(sz); r.font.bold = bold
            r.font.color.rgb = col
    return tb

def pic_w(s, path, x, y, w):  # fixed width, keep aspect
    return s.shapes.add_picture(path, Inches(x), Inches(y), width=Inches(w))
def pic_h(s, path, x, y, h):  # fixed height, keep aspect
    return s.shapes.add_picture(path, Inches(x), Inches(y), height=Inches(h))

def chip(s, x, y, label, col):
    c = box(s, x, y, 0.34, 0.34, fill=col, radius=True)
    return c

# =====================================================================
# SLIDE 1 — TITLE
# =====================================================================
s = slide(NAVY)
box(s, 0, 0, SW, 0.16, fill=CORAL)  # thin top energy strip is fine as a single brand mark
text(s, 0.9, 1.5, 11.7, 1.7,
     [[("레퍼런스 기반 ", 38, LIGHT, True), ("한국어 카드뉴스", 38, CORAL, True)],
      [("레이아웃 인식 ", 38, LIGHT, True), ("& 콘텐츠 치환", 38, AMBER, True)]],
     line_spacing=1.08)
text(s, 0.9, 3.35, 11.5, 1.1,
     [[("레퍼런스 카드의 '형식'을 AI로 정확히 인식해, ", 18, MUTE_D, False)],
      [("같은 포맷에 내 제품 콘텐츠만 바꿔 끼우는 디자인 자동화", 18, MUTE_D, False)]],
     line_spacing=1.15)
# meta row
text(s, 0.9, 5.7, 11.5, 0.5,
     [[("Computer Vision 학기 프로젝트", 15, LIGHT, True),
       ("   ·   YOLOv8 전이학습 + 소규모 데이터 ablation", 15, MUTE_D, False)]])
text(s, 0.9, 6.25, 11.5, 0.4,
     [[("iluv4   ·   2026-06-16   ·   github.com/iluv4/cardnews-cv", 13, MUTE_D, False)]])
# stat trio bottom-right
stats = [("0.85", "mAP@50", MINT), ("24", "실험 (ablation)", AMBER), ("5-fold", "교차검증", CORAL)]
sx = 8.0
for v, lab, col in stats:
    text(s, sx, 4.55, 1.7, 0.6, [[(v, 30, col, True)]])
    text(s, sx, 5.2, 1.7, 0.4, [[(lab, 11, MUTE_D, False)]])
    sx += 1.75

# =====================================================================
# SLIDE 2 — 문제 & 목표
# =====================================================================
s = slide(LIGHT)
text(s, 0.9, 0.55, 11.5, 0.8, [[("문제 & 목표", 32, INK, True)]])
text(s, 0.9, 1.35, 11.5, 0.5,
     [[("잘 만든 카드뉴스의 ", 16, MUTE_L, False), ("'형식'을 재활용", 16, CORAL, True),
       (" 하고 싶다 — 매번 새로 디자인하지 않고", 16, MUTE_L, False)]])
# left: problem  / right: goal
box(s, 0.9, 2.15, 5.55, 4.6, fill=TINT, radius=True)
text(s, 1.25, 2.45, 4.9, 0.5, [[("기존 방식의 한계", 18, INK, True)]])
for i, t in enumerate([
    "카드뉴스 한 장마다 레이아웃을 수작업으로 배치",
    "잘 나온 레퍼런스가 있어도 그대로 못 가져옴",
    "LLM 배치는 위치가 들쭉날쭉, 디자인 일관성 부족"]):
    cy = 3.15 + i*0.95
    chip(s, 1.25, cy, "", CORAL)
    text(s, 1.75, cy-0.04, 4.45, 0.85, [[(t, 14, INK, False)]], anchor=MSO_ANCHOR.TOP)

box(s, 6.85, 2.15, 5.6, 4.6, fill=NAVY, radius=True, shadow=True)
text(s, 7.2, 2.45, 4.95, 0.5, [[("이 프로젝트의 목표", 18, LIGHT, True)]])
for i, (h, d, col) in enumerate([
    ("① 레이아웃 정확 인식", "레퍼런스에서 제목·본문·로고·장식 위치를 검출", MINT),
    ("② 포맷 유지 + 콘텐츠 치환", "같은 구도에 '내 제품' 텍스트/이미지만 교체", AMBER),
    ("→ 디자인 자동화", "레퍼런스를 템플릿처럼 재사용", CORAL)]):
    cy = 3.2 + i*1.15
    text(s, 7.2, cy, 4.95, 0.4, [[(h, 15, col, True)]])
    text(s, 7.2, cy+0.42, 4.95, 0.6, [[(d, 12.5, MUTE_D, False)]], line_spacing=1.1)

# =====================================================================
# SLIDE 3 — 왜 CV 과제인가 (접근)
# =====================================================================
s = slide(LIGHT)
text(s, 0.9, 0.55, 11.5, 0.8, [[("접근: 핵심은 ", 32, INK, True), ("딥러닝 검출 모델 학습", 32, CORAL, True)]])
text(s, 0.9, 1.4, 11.5, 0.5,
     [[("'레이아웃 인식' = 4개 요소를 검출·분류하는 object detection 문제로 정식화", 16, MUTE_L, False)]])
# 4 class cards
classes = [("제목 / title", "큰 헤드라인 텍스트", CORAL),
           ("본문 / body", "설명·캡션 텍스트", MINT),
           ("로고 / logo", "브랜드 마크", AMBER),
           ("장식 / underlay", "배경 박스·강조 도형", RGBColor(0x6C,0x8C,0xFF))]
cw = 2.85; gap = 0.18; x0 = 0.9
for i,(h,d,col) in enumerate(classes):
    x = x0 + i*(cw+gap)
    box(s, x, 2.2, cw, 1.7, fill=CARDBG, radius=True, shadow=True)
    box(s, x+0.28, 2.5, 0.16, 0.16, fill=col, radius=True)
    text(s, x+0.55, 2.42, cw-0.7, 0.4, [[(h, 14.5, INK, True)]])
    text(s, x+0.28, 3.0, cw-0.5, 0.7, [[(d, 12, MUTE_L, False)]], line_spacing=1.1)
# CV techniques row
text(s, 0.9, 4.3, 11.5, 0.5, [[("이 과정에서 다루는 CV 기법", 17, INK, True)]])
techs = ["전이학습 (pretrained → fine-tune)", "객체 검출 / 바운딩박스 회귀",
         "데이터 증강 (augmentation)", "mAP 평가지표", "k-fold 교차검증", "하이퍼파라미터 ablation"]
tw = 3.75; tgap = 0.2
for i,t in enumerate(techs):
    r = i//3; c = i%3
    x = 0.9 + c*(tw+tgap); y = 4.95 + r*0.85
    box(s, x, y, tw, 0.65, fill=TINT, radius=True)
    text(s, x+0.25, y, tw-0.4, 0.65, [[(t, 12.5, INK, False)]], anchor=MSO_ANCHOR.MIDDLE)

# =====================================================================
# SLIDE 4 — 데이터 & 라벨링
# =====================================================================
s = slide(LIGHT)
text(s, 0.9, 0.55, 7.0, 0.8, [[("데이터 & 라벨링", 32, INK, True)]])
text(s, 0.9, 1.4, 7.2, 1.0,
     [[("한국어 인스타 카드뉴스를 직접 수집하고, ", 15, MUTE_L, False)],
      [("EasyOCR로 ", 15, MUTE_L, False), ("의사 라벨(pseudo-label)", 15, CORAL, True),
       ("을 만들어 YOLO 형식으로 변환", 15, MUTE_L, False)]], line_spacing=1.15)
# stat callouts
for i,(v,lab,col) in enumerate([("109","수집한 카드 이미지", MINT),
                                ("4","검출 클래스", AMBER),
                                ("925","자동 생성 박스 (제목284·본문641)", CORAL)]):
    y = 2.7 + i*1.25
    text(s, 0.9, y, 2.0, 0.7, [[(v, 36, col, True)]])
    text(s, 2.7, y+0.12, 4.4, 0.8, [[(lab, 13, INK, False)]], anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.05)
# right: annotated reference card
box(s, 8.0, 1.35, 4.45, 5.55, fill=CARDBG, radius=True, shadow=True)
pic_h(s, P("label_preview","img_022.png"), 8.45, 1.75, 4.4)
text(s, 8.0, 6.35, 4.45, 0.5,
     [[("자동 라벨 예시 — 빨강=제목, 초록=본문", 11.5, MUTE_L, False)]], align=PP_ALIGN.CENTER)

# =====================================================================
# SLIDE 5 — 결과: 작동한다 (best model)
# =====================================================================
s = slide(NAVY)
text(s, 0.9, 0.55, 11.5, 0.8, [[("결과 — ", 32, LIGHT, True), ("실제로 인식한다", 32, MINT, True)]])
text(s, 0.9, 1.4, 11.5, 0.5,
     [[("최고 모델 ", 15, MUTE_D, False), ("e15_long300_card", 15, AMBER, True),
       (" (YOLOv8n, 300 epochs)", 15, MUTE_D, False)]])
# big metrics
mets = [("0.854","mAP@50", MINT), ("0.718","mAP@50-95", CORAL),
        ("0.785","Precision", AMBER), ("0.842","Recall", RGBColor(0x6C,0x8C,0xFF))]
for i,(v,lab,col) in enumerate(mets):
    y = 2.35 + i*1.05
    text(s, 0.9, y, 2.2, 0.65, [[(v, 30, col, True)]])
    text(s, 0.9, y+0.6, 2.5, 0.3, [[(lab, 12, MUTE_D, False)]])
# right: prediction grid
box(s, 5.0, 1.95, 7.45, 4.95, fill=PANEL, radius=True, shadow=True)
pic_h(s, P("results","e15_long300_card","val_batch0_pred.jpg"), 5.3, 2.2, 4.45)
text(s, 9.9, 2.4, 2.4, 4.2,
     [[("검증셋 예측", 15, LIGHT, True)],
      [("", 6, LIGHT, False)],
      [("처음 보는 카드에서도 제목·본문 영역을 정확히 박스로 검출", 13, MUTE_D, False)]],
     anchor=MSO_ANCHOR.TOP, line_spacing=1.2)

# =====================================================================
# SLIDE 6 — Ablation 핵심 발견 (native chart)
# =====================================================================
s = slide(LIGHT)
text(s, 0.9, 0.55, 11.5, 0.8, [[("Ablation — ", 32, INK, True), ("무엇이 성능을 올렸나", 32, CORAL, True)]])
text(s, 0.9, 1.4, 11.5, 0.5,
     [[("24개 실험을 비교 (mAP@50-95 기준)", 15, MUTE_L, False)]])
# chart
cd = CategoryChartData()
cd.categories = ["Scratch\n(전이학습X)", "Baseline\nYOLOv8n", "+ Card Aug", "YOLOv8s", "Best\n(e15·300ep)"]
cd.add_series("mAP@50-95", (0.502, 0.662, 0.696, 0.705, 0.718))
gf = s.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED,
                        Inches(0.9), Inches(2.15), Inches(7.0), Inches(4.7), cd)
ch = gf.chart
ch.has_legend = False; ch.has_title = False
plot = ch.plots[0]; plot.has_data_labels = True
plot.gap_width = 60
dl = plot.data_labels; dl.number_format = '0.00'; dl.number_format_is_linked = False
dl.position = XL_LABEL_POSITION.OUTSIDE_END
dl.font.size = Pt(12); dl.font.bold = True; dl.font.color.rgb = INK
ser = plot.series[0]
from pptx.oxml.ns import qn
# color each point
colors = ['B8C0D8','6C8CFF','16C08E','F5A623','FF5C7A']
for idx, hexc in enumerate(colors):
    pt = ser.points[idx]
    pt.format.fill.solid(); pt.format.fill.fore_color.rgb = RGBColor.from_string(hexc)
cat = ch.category_axis; cat.tick_labels.font.size = Pt(10); cat.tick_labels.font.color.rgb = MUTE_L
cat.format.line.color.rgb = RGBColor(0xD0,0xD6,0xE6)
val = ch.value_axis; val.minimum_scale = 0.0; val.maximum_scale = 0.8
val.tick_labels.font.size = Pt(10); val.tick_labels.font.color.rgb = MUTE_L
val.has_major_gridlines = True
val.major_gridlines.format.line.color.rgb = RGBColor(0xEC,0xEF,0xF7)
val.major_gridlines.format.line.width = Pt(0.5)
# findings panel
box(s, 8.2, 2.15, 4.25, 4.7, fill=NAVY, radius=True, shadow=True)
text(s, 8.55, 2.45, 3.6, 0.4, [[("핵심 발견", 17, LIGHT, True)]])
finds = [("전이학습이 결정적", "scratch 0.50 → fine-tune 0.72", MINT),
         ("카드 특화 증강 효과", "card aug로 +0.03", AMBER),
         ("YOLOv8s ≥ n", "조금 더 크면 소폭 향상", RGBColor(0x6C,0x8C,0xFF)),
         ("길게 학습 + 증강", "300ep가 최고 성능", CORAL)]
for i,(h,d,col) in enumerate(finds):
    y = 3.1 + i*0.92
    text(s, 8.55, y, 3.6, 0.35, [[("• ", 13, col, True),(h, 13.5, LIGHT, True)]])
    text(s, 8.75, y+0.36, 3.45, 0.4, [[(d, 11.5, MUTE_D, False)]])

# =====================================================================
# SLIDE 7 — 신뢰성: 학습곡선 + k-fold
# =====================================================================
s = slide(LIGHT)
text(s, 0.9, 0.55, 11.5, 0.8, [[("신뢰성 — 학습 곡선 & 교차검증", 30, INK, True)]])
text(s, 0.9, 1.35, 11.5, 0.5,
     [[("과적합 없이 수렴, 그리고 5-fold로 일반화 성능 확인", 15, MUTE_L, False)]])
box(s, 0.9, 2.1, 7.6, 4.0, fill=CARDBG, radius=True, shadow=True)
pic_w(s, P("results","e15_long300_card","results.png"), 1.15, 2.5, 7.1)
text(s, 0.9, 6.2, 7.6, 0.4, [[("학습 곡선 (e15, 300 epochs) — loss 감소, mAP 안정 수렴", 11.5, MUTE_L, False)]], align=PP_ALIGN.CENTER)
# kfold panel
box(s, 8.75, 2.1, 3.7, 4.0, fill=TINT, radius=True)
text(s, 9.1, 2.4, 3.0, 0.4, [[("5-fold 교차검증", 16, INK, True)]])
text(s, 9.1, 3.0, 3.0, 0.9, [[("0.61", 40, MINT, True)],[("평균 mAP@50-95", 12, MUTE_L, False)]])
text(s, 9.1, 4.2, 3.1, 1.7,
     [[("표준편차 ±0.05", 13, INK, True)],
      [("", 5, INK, False)],
      [("폴드별: 0.65 / 0.57 / 0.68 / 0.60 / 0.55", 11.5, MUTE_L, False)],
      [("", 4, INK, False)],
      [("→ 109장 소규모치고 안정적", 12, INK, True)]], line_spacing=1.15)

# =====================================================================
# SLIDE 8 — 현재 상태 & 다음 단계
# =====================================================================
s = slide(NAVY)
text(s, 0.9, 0.55, 11.5, 0.8, [[("현재 상태 & 다음 단계", 32, LIGHT, True)]])
# done
box(s, 0.9, 1.7, 5.6, 5.0, fill=PANEL, radius=True, shadow=True)
text(s, 1.25, 2.0, 4.9, 0.5, [[("✓ 완료", 18, MINT, True)]])
for i,t in enumerate([
    "한국어 카드 109장 수집 + 자동 라벨링",
    "YOLOv8 전이학습 — 24개 실험 ablation",
    "최고 모델 mAP@50 0.85 / mAP@50-95 0.72",
    "5-fold 교차검증으로 일반화 확인",
    "재현 가능한 코드 + best.pt 공개"]):
    y = 2.65 + i*0.8
    chip(s, 1.25, y, "", MINT)
    text(s, 1.75, y-0.05, 4.5, 0.75, [[(t, 13, LIGHT, False)]], line_spacing=1.05)
# next
box(s, 6.85, 1.7, 5.6, 5.0, fill=PANEL, radius=True, shadow=True)
text(s, 7.2, 2.0, 4.9, 0.5, [[("→ 다음", 18, AMBER, True)]])
for i,(t,d) in enumerate([
    ("콘텐츠 자동 치환", "인식한 레이아웃에 내 제품 텍스트/이미지 배치"),
    ("로고·장식 수동 라벨", "현재 제목·본문만 → 4클래스 완성"),
    ("데이터 확장", "109장 → 1,000장+ 로 정확도 향상"),
    ("앱 통합", "레퍼런스 업로드 → 자동 카드 생성")]):
    y = 2.65 + i*1.0
    chip(s, 7.2, y, "", AMBER)
    text(s, 7.7, y-0.05, 4.5, 0.45, [[(t, 13.5, LIGHT, True)]])
    text(s, 7.7, y+0.4, 4.5, 0.5, [[(d, 11.5, MUTE_D, False)]], line_spacing=1.05)

prs.save(P("cardnews_layout_detector.pptx"))
print("saved:", P("cardnews_layout_detector.pptx"))
