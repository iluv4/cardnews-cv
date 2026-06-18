# -*- coding: utf-8 -*-
"""Build the cardnews-cv presentation deck (python-pptx).
  py -3 slides/make_deck.py   ->  slides/cardnews-cv.pptx
Dark "forest + pink" palette to match the product's rendered cards.
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
def P(*p): return os.path.join(ROOT, *p)

# palette
BG    = RGBColor(0x0F, 0x1F, 0x18)
PANEL = RGBColor(0x18, 0x2D, 0x23)
PANEL2= RGBColor(0x1E, 0x38, 0x2C)
GREEN = RGBColor(0x2A, 0x7A, 0x52)
PINK  = RGBColor(0xF7, 0x88, 0x9A)
MINT  = RGBColor(0x5C, 0xC8, 0xA0)
TEXT  = RGBColor(0xEA, 0xF2, 0xEC)
MUTED = RGBColor(0x9D, 0xB3, 0xA7)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LINE  = RGBColor(0x2E, 0x4A, 0x3C)
FONT  = "Malgun Gothic"          # Korean-safe on Windows

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
SW, SH = prs.slide_width, prs.slide_height


def slide(bg=BG):
    s = prs.slides.add_slide(BLANK)
    r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    r.fill.solid(); r.fill.fore_color.rgb = bg; r.line.fill.background()
    r.shadow.inherit = False
    s.shapes._spTree.remove(r._element); s.shapes._spTree.insert(2, r._element)
    return s


def box(s, l, t, w, h, anchor=MSO_ANCHOR.TOP):
    tb = s.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Pt(0)
    tf.margin_top = tf.margin_bottom = Pt(0)
    return tf


def para(tf, text, size, color=TEXT, bold=False, italic=False, first=False,
         align=PP_ALIGN.LEFT, space_after=6, bullet=False, font=FONT):
    p = tf.paragraphs[0] if first and not tf.paragraphs[0].runs else tf.add_paragraph()
    p.alignment = align; p.space_after = Pt(space_after); p.space_before = Pt(0)
    r = p.add_run(); r.text = text
    f = r.font; f.size = Pt(size); f.bold = bold; f.italic = italic
    f.name = font; f.color.rgb = color
    if bullet:
        _bullet(p)
    return p


def _bullet(p, char="·", color=PINK):
    pPr = p._pPr if p._pPr is not None else p.get_or_add_pPr()
    bu = pPr.makeelement(qn('a:buChar'), {'char': char})
    pPr.append(bu)


def rect(s, l, t, w, h, fill=PANEL, line=None, rounded=False, line_w=1.0):
    shp = s.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE,
        Inches(l), Inches(t), Inches(w), Inches(h))
    shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line; shp.line.width = Pt(line_w)
    shp.shadow.inherit = False
    return shp


def circle_num(s, l, t, n, d=0.46, fill=PINK, fg=RGBColor(0x2A,0x0D,0x14)):
    c = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(l), Inches(t), Inches(d), Inches(d))
    c.fill.solid(); c.fill.fore_color.rgb = fill; c.line.fill.background()
    c.shadow.inherit = False
    tf = c.text_frame; tf.margin_top = tf.margin_bottom = Pt(0)
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = str(n)
    r.font.size = Pt(16); r.font.bold = True; r.font.name = FONT; r.font.color.rgb = fg
    return c


def pic_framed(s, path, l, t, w, h):
    """rounded mat behind a contained image."""
    rect(s, l-0.08, t-0.08, w+0.16, h+0.16, fill=PANEL2, line=LINE, rounded=True)
    from PIL import Image
    iw, ih = Image.open(path).size
    ar = iw/ih; box_ar = w/h
    if ar > box_ar:
        nw = w; nh = w/ar
    else:
        nh = h; nw = h*ar
    s.shapes.add_picture(path, Inches(l+(w-nw)/2), Inches(t+(h-nh)/2), Inches(nw), Inches(nh))


def kicker(s, text, l=0.9, t=0.55):
    para(box(s, l, t, 8, 0.4), text, 13, MINT, bold=True, first=True)

def title(s, text, l=0.9, t=0.92, w=11.5, size=34):
    para(box(s, l, t, w, 1.0), text, size, WHITE, bold=True, first=True)


# ---------------- Slide 1: Title ----------------
s = slide()
para(box(s, 0.9, 2.35, 11.5, 1.4), "카드뉴스-CV", 60, WHITE, bold=True, first=True)
para(box(s, 0.95, 3.7, 11.5, 0.8), "레퍼런스 검색 기반 한국어 카드뉴스 생성 엔진", 24, PINK, bold=True, first=True)
tf = box(s, 0.95, 4.55, 11.5, 1.2)
para(tf, "분석·검색 중심 ML  ·  결정적 렌더링 (GAN 픽셀 생성 X)", 16, TEXT, first=True)
para(tf, "검색 → 선택 → 레이아웃 복사 → 디자이너급 자동 채움", 16, MUTED)
para(box(s, 0.95, 6.6, 11.5, 0.5), "2026-06-18   ·   github.com/iluv4/cardnews-cv", 12, MUTED, first=True)

# ---------------- Slide 2: Problem & insight ----------------
s = slide()
kicker(s, "문제 & 인사이트")
title(s, "‘자동 생성’은 왜 ‘AI같아’ 보이는가")
cards = [
    (0.9, "GAN · 디퓨전의 한계", PINK, [
        "한글 텍스트를 또렷하게 못 그림",
        "결과가 불안정하고 통제 어려움",
        "흔한 플랫 그라데이션 = ‘AI-tic’",
    ]),
    (6.95, "우리의 접근", MINT, [
        "ML은 분석(검출기) + 검색(레퍼런스)에",
        "픽셀은 결정적 렌더링으로 또렷하게",
        "실제 디자인 레퍼런스를 복사 → 디자이너급",
    ]),
]
for x, head, col, items in cards:
    rect(s, x, 2.1, 5.45, 4.4, fill=PANEL, line=LINE, rounded=True)
    para(box(s, x+0.4, 2.45, 4.7, 0.6), head, 22, col, bold=True, first=True)
    tf = box(s, x+0.4, 3.3, 4.7, 3.0)
    for i, it in enumerate(items):
        para(tf, it, 16, TEXT, first=(i == 0), space_after=12)

# ---------------- Slide 3: Pipeline ----------------
s = slide()
kicker(s, "시스템 개요")
title(s, "파이프라인 — 손으로 짠 한 스타일이 아니다")
steps = [
    ("검출기", "레이아웃 이해\ntitle·body·logo·underlay"),
    ("레퍼런스 검색", "의미·색·구조로\n원하는 디자인 찾기"),
    ("레이아웃 복사", "선택한 레퍼런스의\n실제 구조를 추출"),
    ("v2 렌더", "그 구조에 사용자 텍스트를\n디자이너급으로 채움"),
]
x0, w, gap = 0.9, 2.85, 0.27
for i, (h, d) in enumerate(steps):
    x = x0 + i*(w+gap)
    rect(s, x, 2.35, w, 3.0, fill=PANEL, line=LINE, rounded=True)
    circle_num(s, x+0.35, 2.7, i+1)
    para(box(s, x+0.32, 3.45, w-0.6, 0.6), h, 19, WHITE, bold=True, first=True)
    tf = box(s, x+0.32, 4.1, w-0.6, 1.1)
    for j, ln in enumerate(d.split("\n")):
        para(tf, ln, 13, MUTED, first=(j == 0), space_after=2)
    if i < 3:
        para(box(s, x+w-0.02, 3.5, 0.3, 0.5), "→", 20, PINK, first=True, align=PP_ALIGN.CENTER)
para(box(s, 0.9, 5.7, 11.5, 0.6), "→ 엔진 수 = 데이터셋 수: 레퍼런스마다 레이아웃 템플릿 (손코딩 아님)",
     15, MINT, italic=True, first=True)

# ---------------- Slide 4: Detector ----------------
s = slide()
kicker(s, "기반 기술")
title(s, "레이아웃 검출기 — 완료")
stats = [
    ("0.718", "mAP@50-95 (best)"),
    ("0.854", "mAP@50"),
    ("0.611±0.049", "5-fold CV"),
    ("0.718 vs 0.502", "전이학습 vs 스크래치"),
]
x0, w, gap = 0.9, 2.85, 0.27
for i, (big, lab) in enumerate(stats):
    x = x0 + i*(w+gap)
    rect(s, x, 2.5, w, 2.3, fill=PANEL, line=LINE, rounded=True)
    para(box(s, x+0.2, 2.95, w-0.4, 0.9, MSO_ANCHOR.MIDDLE), big,
         30 if len(big) < 7 else 20, PINK, bold=True, first=True, align=PP_ALIGN.CENTER)
    para(box(s, x+0.2, 3.95, w-0.4, 0.7), lab, 13, MUTED, first=True, align=PP_ALIGN.CENTER)
para(box(s, 0.9, 5.5, 11.5, 0.6),
     "687장 (109 시드 + 578 신규, 66덱) · YOLO · imgsz 1280 · 전이학습 결정적",
     15, TEXT, first=True)

# ---------------- Slide 5: Reference library + search ----------------
s = slide()
kicker(s, "제품 · ML 핵심")
title(s, "레퍼런스 라이브러리 + 검색")
tf = box(s, 0.9, 2.2, 6.7, 4.5)
items = [
    "796장 인덱싱 — 팔레트·밝기·엣지·다크 + 파일메타",
    "8개 시각 아키타입으로 클러스터",
    "검색: 텍스트(주제)·색·필터(다크/표지/클러스터)·유사도",
    "torch 없이 로컬 동작 — 지금 바로 검색 가능",
    "CLIP 의미검색 + 검출기 태깅은 RunPod 훅으로 대기",
]
for i, it in enumerate(items):
    para(tf, it, 17, TEXT, first=(i == 0), space_after=14, bullet=True)
pic_framed(s, P("service", "library", "contact_sheet.png"), 8.05, 1.95, 4.4, 4.9)
para(box(s, 8.05, 6.85, 4.4, 0.4), "레퍼런스 코퍼스에서 추출한 레이아웃 템플릿", 11, MUTED, first=True, align=PP_ALIGN.CENTER)

# ---------------- Slide 6: Engine v2 quality ----------------
s = slide()
kicker(s, "렌더 품질")
title(s, "실제 레퍼런스 디자인 언어를 맞추다")
tf = box(s, 0.9, 2.3, 5.0, 4.2)
for i, it in enumerate([
    "다크 배경 + 스캐터 데코 레이어",
    "2톤 글로우 타이틀 + 키워드 강조 부제",
    "화이트 라운드 패널 + 액센트 체크리스트",
    "브랜드 · 페이지 · 마스코트 슬롯",
]):
    para(tf, it, 17, TEXT, first=(i == 0), space_after=14, bullet=True)
para(box(s, 0.9, 6.2, 5.0, 0.6), "‘AI같음’ 격차를 닫은 결정적 렌더", 14, MINT, italic=True, first=True)
pic_framed(s, P("service", "out", "deck_00.png"), 6.2, 2.0, 3.4, 3.4)
pic_framed(s, P("service", "out", "deck_01.png"), 9.75, 2.0, 3.4, 3.4)

# ---------------- Slide 7: Layout copy + refill ----------------
s = slide()
kicker(s, "진짜 엔진")
title(s, "레이아웃 복사 → 사용자 텍스트 리필")
tf = box(s, 0.9, 2.3, 5.6, 4.2)
for i, it in enumerate([
    "선택 레퍼런스의 title·body 블록 + 팔레트 추출",
    "사용자 텍스트로 리필 (테마는 레퍼런스 색에서 합성)",
    "텍스트 길이가 달라도 패널 침범 방지 (clamp)",
    "템플릿 없으면 테마 폴백 — RunPod로 전 레퍼런스 확장",
]):
    para(tf, it, 16, TEXT, first=(i == 0), space_after=13, bullet=True)
pic_framed(s, P("service", "out", "from_ref_template.png"), 6.7, 1.95, 3.1, 3.9)
pic_framed(s, P("gen_output", "refill", "refill_list_img_001.png"), 10.0, 2.55, 3.0, 1.9)
para(box(s, 10.0, 4.6, 3.0, 0.4), "원본 | 리필 비교", 11, MUTED, first=True, align=PP_ALIGN.CENTER)

# ---------------- Slide 8: Service + UI ----------------
s = slide()
kicker(s, "서비스")
title(s, "검색 → 선택 → 생성 → 내보내기")
eps = [
    ("GET", "/api/search · /similar · /clusters", "레퍼런스 검색"),
    ("GET", "/api/reference/{id}", "썸네일 제공"),
    ("POST", "/api/generate_from_ref", "선택 레퍼런스 레이아웃 복사"),
    ("POST", "/api/generate · /api/deck", "단일 카드 PNG · 덱 ZIP"),
]
y = 2.3
for m, path, desc in eps:
    rect(s, 0.9, y, 11.5, 0.92, fill=PANEL, line=LINE, rounded=True)
    b = rect(s, 1.15, y+0.22, 1.0, 0.48, fill=GREEN, rounded=True)
    tf = b.text_frame; tf.margin_top=Pt(0); tf.margin_bottom=Pt(0)
    p = tf.paragraphs[0]; p.alignment=PP_ALIGN.CENTER
    rr=p.add_run(); rr.text=m; rr.font.size=Pt(13); rr.font.bold=True; rr.font.name=FONT; rr.font.color.rgb=WHITE
    para(box(s, 2.4, y+0.13, 6.0, 0.6), path, 16, TEXT, bold=True, first=True, font="Consolas")
    para(box(s, 8.4, y+0.16, 3.8, 0.6), desc, 13, MUTED, first=True)
    y += 1.06
para(box(s, 0.9, 6.75, 11.5, 0.5), "웹 UI · WCAG 2.1 AA 통과 — 키보드만으로 전 기능 사용 가능",
     14, MINT, italic=True, first=True)

# ---------------- Slide 9: Demo ----------------
s = slide()
kicker(s, "라이브 데모")
title(s, "직접 보여드립니다")
steps = [
    ".\\run_demo.ps1  →  http://127.0.0.1:8000/",
    "‘smart farm’ 검색 → 초록 ‘레이아웃’ 배지 카드 선택",
    "제목·체크리스트 채우고 [생성] 클릭",
    "(옵션) 덱 전체를 ZIP으로 내보내기",
]
y = 2.4
for i, st in enumerate(steps):
    circle_num(s, 0.95, y, i+1)
    para(box(s, 1.7, y-0.02, 11.0, 0.6, MSO_ANCHOR.MIDDLE), st, 18, TEXT, first=True)
    y += 0.95
para(box(s, 0.95, 6.5, 11.5, 0.5),
     "라이브가 막히면 폴백: service/out/ · gen_output/refill/ 의 미리 렌더된 결과",
     13, MUTED, italic=True, first=True)

# ---------------- Slide 10: Results summary ----------------
s = slide()
kicker(s, "지금까지")
title(s, "무엇이 동작하나")
tf = box(s, 0.9, 2.25, 7.4, 4.5)
for i, it in enumerate([
    "검출기 mAP 0.718 · reflib 검색 796장 · 8 아키타입 클러스터",
    "엔진 v2 — 디자이너급 렌더 (레퍼런스 디자인 언어)",
    "레이아웃 복사 + 리필 — 95개 템플릿, 서비스에 연결",
    "FastAPI + 웹 UI + WCAG 2.1 AA",
    "전부 로컬 동작 · 커밋·푸시 완료 (main)",
]):
    para(tf, it, 16, TEXT, first=(i == 0), space_after=13, bullet=True)
rect(s, 8.7, 2.5, 3.7, 3.3, fill=PANEL, line=LINE, rounded=True)
para(box(s, 8.9, 3.0, 3.3, 1.2, MSO_ANCHOR.MIDDLE), "4 / 5", 60, PINK, bold=True, first=True, align=PP_ALIGN.CENTER)
para(box(s, 8.9, 4.4, 3.3, 1.0), "로드맵 항목 완료\n(남은 1개 = RunPod 패스)", 15, MUTED, first=True, align=PP_ALIGN.CENTER)

# ---------------- Slide 11: Next ----------------
s = slide()
kicker(s, "다음 단계")
title(s, "RunPod 한 패스로 완성")
tf = box(s, 0.9, 2.4, 11.4, 3.0)
for i, it in enumerate([
    "검출기를 687장 전체에 실행 → 템플릿 확장 (테마 폴백이 사라짐)",
    "CLIP 임베딩 빌드 → 의미검색 ON (lexical → 의미 기반 자동 승격)",
    "명령 한 줄: reflib/embed_clip.py --build · reflib/tag_layout.py",
]):
    para(tf, it, 18, TEXT, first=(i == 0), space_after=16, bullet=True)
rect(s, 0.9, 5.5, 11.5, 1.1, fill=PANEL, line=LINE, rounded=True)
para(box(s, 1.2, 5.72, 11.0, 0.7, MSO_ANCHOR.MIDDLE),
     "→ 검색한 어떤 레퍼런스든 실제 레이아웃을 복사하게 됨", 18, MINT, bold=True, first=True)

# ---------------- Slide 12: Closing ----------------
s = slide()
para(box(s, 0.9, 2.9, 11.5, 1.2), "감사합니다", 48, WHITE, bold=True, first=True)
para(box(s, 0.95, 4.1, 11.5, 0.6), "카드뉴스-CV — 분석·검색 중심 한국어 카드뉴스 생성", 18, PINK, first=True)
para(box(s, 0.95, 4.8, 11.5, 0.5), "github.com/iluv4/cardnews-cv  ·  데모: .\\run_demo.ps1", 14, MUTED, first=True)

out = P("slides", "cardnews-cv.pptx")
prs.save(out)
print("saved", out, "-", len(prs.slides.__iter__.__self__._sldIdLst), "slides")
